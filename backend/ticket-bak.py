from fastapi import APIRouter, File, UploadFile, HTTPException, BackgroundTasks, Query
from fastapi.responses import StreamingResponse, FileResponse
from api.model.models import (
    ResData,
    BatchListResponse,
    BatchInfo,
    BatchDetail,
    BatchStatusResponse,
    UploadResponse,
    DashboardStats,
    SystemStatus,
    PassengerInfo,
    create_success_response,
    calculate_progress_percentage
)
from api.utils.pdf_converter import get_converter
from api.utils.batch_manager import get_batch_manager
from api.utils.agency_manager import get_agency_manager
from api.utils.file_handler import get_file_handler
import pandas as pd
from dataclasses import dataclass
from typing import Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from datetime import datetime
import traceback


router = APIRouter(prefix="/api/v1/ticket")


# ============================================================================
# DATA CLASSES (Keep your existing Ticket dataclass)
# ============================================================================

@dataclass
class Ticket:
    no: int
    rsvn_cfmd: str
    ticket_type: str
    pax_name: str
    emd1: str
    pnr: str
    travel_agency: str
    ptn1_dep: str
    ptn1_dep_date: str
    ptn1_dep_time: str
    ptn1_arr: str
    ptn1_arr_date: str
    ptn1_arr_time: str
    ptn2_dep: Optional[str] = None
    ptn2_dep_date: Optional[str] = None
    ptn2_dep_time: Optional[str] = None
    ptn2_arr: Optional[str] = None
    ptn2_arr_date: Optional[str] = None
    ptn2_arr_time: Optional[str] = None


# ============================================================================
# HELPER FUNCTIONS (Keep your existing PPTX functions)
# ============================================================================

def replace_text_in_paragraph(paragraph, replacements):
    """Replace text in a paragraph by combining all runs first"""
    full_text = paragraph.text
    original_text = full_text
    
    for old_text, new_text in replacements.items():
        if old_text in full_text:
            full_text = full_text.replace(old_text, str(new_text))
    
    if full_text != original_text:
        # Preserve the first run's formatting
        if paragraph.runs:
            first_run = paragraph.runs[0]
            # Store formatting
            font_name = first_run.font.name
            font_size = first_run.font.size
            bold = first_run.font.bold
            italic = first_run.font.italic
            
            # Clear all runs
            for run in paragraph.runs:
                run.text = ""
            
            # Set new text with preserved formatting
            first_run.text = full_text
            if font_name:
                first_run.font.name = font_name
            if font_size:
                first_run.font.size = font_size
            if bold is not None:
                first_run.font.bold = bold
            if italic is not None:
                first_run.font.italic = italic
        else:
            paragraph.add_run(full_text)


def replace_text_in_shape(shape, replacements):
    """Replace text in a shape, handling Groups, Tables, and TextFrames recursively"""
    
    # Handle groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            replace_text_in_shape(child_shape, replacements)
        return

    # Handle text frames
    if shape.has_text_frame:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            replace_text_in_paragraph(paragraph, replacements)
    
    # Handle tables
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replace_text_in_paragraph(paragraph, replacements)


def generate_ppt_from_template(template_path, ticket, agency=None):
    """Generate PowerPoint from template with ticket data and agency info"""
    prs = Presentation(template_path)
    
    current_date = datetime.now().strftime("%Y-%m-%d")
    
    # Basic ticket replacements
    replacements = {
        '{{date_now}}': current_date,
        '{{PAX_name}}': ticket.pax_name,
        '{{PNR_Reference}}': ticket.pnr,
        '{{PTN1-Dep}}': ticket.ptn1_dep,
        '{{PTN1-Arr}}': ticket.ptn1_arr,
        '{{PTN1_Date}}': ticket.ptn1_dep_date,
        '{{PTN1_Time}}': ticket.ptn1_dep_time,
        '{{PTN1_Date.1}}': ticket.ptn1_arr_date,
        '{{PTN1_Time.1}}': ticket.ptn1_arr_time,
    }
    
    # Add PTN2 replacements even if empty (to clear placeholders)
    replacements.update({
        '{{PTN2-Dep}}': ticket.ptn2_dep if ticket.ptn2_dep else '',
        '{{PTN2-Arr}}': ticket.ptn2_arr if ticket.ptn2_arr else '',
        '{{PTN2_Date}}': ticket.ptn2_dep_date if ticket.ptn2_dep_date else '',
        '{{PTN2_Time}}': ticket.ptn2_dep_time if ticket.ptn2_dep_time else '',
        '{{PTN2_Date.1}}': ticket.ptn2_arr_date if ticket.ptn2_arr_date else '',
        '{{PTN2_Time.1}}': ticket.ptn2_arr_time if ticket.ptn2_arr_time else '',
    })
    
    # Add agency information if found
    if agency:
        replacements.update({
            '{{agency_name}}': agency.get('agency_name', ''),
            '{{agency_owner}}': agency.get('agency_owner', ''),
            '{{agency_address}}': agency.get('agency_address', ''),
            '{{agency_email}}': agency.get('email', ''),
            '{{agency_phone}}': agency.get('telephone', '')
        })
        print(f"  ℹ️  Using agency: {agency.get('agency_name')}")
    else:
        # Use defaults if agency not found
        replacements.update({
            '{{agency_name}}': ticket.travel_agency if ticket.travel_agency else '',
            '{{agency_owner}}': '',
            '{{agency_address}}': '',
            '{{agency_email}}': '',
            '{{agency_phone}}': ''
        })
        if ticket.travel_agency:
            print(f"  ⚠️  Agency not found in database: {ticket.travel_agency}")
    
    # Debug: Print all placeholders found in template (optional - can comment out)
    # all_text = []
    # for slide in prs.slides:
    #     for shape in slide.shapes:
    #         if hasattr(shape, "text"):
    #             all_text.append(shape.text)
    # print(f"Template placeholders: {[t for t in all_text if '{{' in t]}")
    
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
    
    return prs


def parse_excel_to_tickets(df: pd.DataFrame) -> list[Ticket]:
    """Parse Excel DataFrame to Ticket objects"""
    df.columns = df.columns.str.replace(' ', '_')
    df = df.fillna("")
    
    tickets = []
    for index, row in df.iterrows():
        raw_emd = row.get('emd1_(Extra_RQ)', '')
        clean_emd = str(int(raw_emd)) if raw_emd != "" else ""

        ticket = Ticket(
            no=row['NO'],
            rsvn_cfmd=row['RSVN_cfmd'],
            ticket_type=row['ADT/LBR/CHD/INF'],
            pax_name=row['PAX_name'],
            emd1=clean_emd, 
            pnr=row['PNR_Reference'],
            travel_agency=row.get('Travel_Agency', ''),
            ptn1_dep=row['PTN1-Dep'],
            ptn1_dep_date=row['PTN1_Date'],
            ptn1_dep_time=row['PTN1_Time'],
            ptn1_arr=row['PTN1-Arr'],
            ptn1_arr_date=row['PTN1_Date.1'],  
            ptn1_arr_time=row['PTN1_Time.1'], 
            ptn2_dep=row.get('PTN2-Dep', ''),
            ptn2_dep_date=row.get('PTN2_Date', ''),
            ptn2_dep_time=row.get('PTN2_Time', ''),
            ptn2_arr=row.get('PTN2-Arr', ''),
            ptn2_arr_date=row.get('PTN2_Date.1', ''),
            ptn2_arr_time=row.get('PTN2_Time.1', '')
        )
        tickets.append(ticket)
    
    return tickets


# ============================================================================
# TICKET GENERATION LOGIC
# ============================================================================

def generate_tickets_for_batch(batch_id: str, tickets: list[Ticket]):
    """
    Background task: Generate PDFs for all tickets in a batch
    This runs after the upload endpoint returns
    """
    manager = get_batch_manager()
    converter = get_converter()
    file_handler = get_file_handler()
    agency_manager = get_agency_manager()
    
    template_path = Path("templates/ticket_template.pptx")
    batch_dir = manager.get_batch_dir(batch_id)
    
    # Update batch status to processing
    manager.update_batch_status(batch_id, "processing")
    
    print(f"\n{'='*60}")
    print(f"Starting batch generation: {batch_id}")
    print(f"Total tickets: {len(tickets)}")
    print(f"{'='*60}\n")
    
    for idx, ticket in enumerate(tickets, 1):
        print(f"[{idx}/{len(tickets)}] Processing: {ticket.pax_name} ({ticket.pnr})")
        
        try:
            # Look up agency from database
            agency = None
            if ticket.travel_agency:
                agency = agency_manager.find_by_name(ticket.travel_agency)
            
            # Generate PPTX with agency info
            prs = generate_ppt_from_template(template_path, ticket, agency)
            
            # Create safe filename
            safe_pax_name = ticket.pax_name.replace('/', '_').replace('\\', '_')
            safe_pnr = ticket.pnr.replace('/', '_')
            
            pptx_filename = f"ticket_{safe_pax_name}_{safe_pnr}.pptx"
            pptx_path = batch_dir / pptx_filename
            
            # Save PPTX temporarily
            prs.save(str(pptx_path))
            print(f"  ✓ PPTX created: {pptx_filename}")
            
            # Convert PPTX to PDF
            pdf_path = converter.convert_pptx_to_pdf(pptx_path, batch_dir)
            pdf_filename = pdf_path.name
            print(f"  ✓ PDF created: {pdf_filename}")
            
            # Update passenger status to generated
            manager.update_passenger_status(
                batch_id=batch_id,
                pax_name=ticket.pax_name,
                pnr=ticket.pnr,
                status="generated",
                pdf_filename=pdf_filename
            )
            
            # Clean up PPTX file
            try:
                pptx_path.unlink()
                print(f"  ✓ Cleaned up PPTX")
            except:
                pass
            
            print(f"  ✅ Completed: {ticket.pax_name}\n")
            
        except Exception as e:
            error_msg = str(e)
            print(f"  ❌ Failed: {error_msg}\n")
            
            # Update passenger status to failed
            manager.update_passenger_status(
                batch_id=batch_id,
                pax_name=ticket.pax_name,
                pnr=ticket.pnr,
                status="failed",
                error=error_msg
            )
    
    # Mark batch as completed
    manager.update_batch_status(batch_id, "completed")
    
    # Get final counts
    batch_info = manager.get_batch_info(batch_id)
    print(f"\n{'='*60}")
    print(f"Batch {batch_id} completed!")
    print(f"Generated: {batch_info['generated']}")
    print(f"Failed: {batch_info['failed']}")
    print(f"{'='*60}\n")


# ============================================================================
# API ENDPOINTS
# ============================================================================

@router.post("/upload", response_model=UploadResponse)
async def upload_excel(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...)
):
    """
    Upload Excel file and create batch
    Returns immediately with batch_id
    Ticket generation happens in background
    """
    # Validate file type
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(
            status_code=400, 
            detail="File must be an Excel file (.xlsx or .xls)"
        )
    
    # Check LibreOffice availability
    converter = get_converter()
    if not converter.is_available():
        raise HTTPException(
            status_code=500,
            detail="LibreOffice not found. Please install LibreOffice from https://www.libreoffice.org/download/"
        )
    
    try:
        # Parse Excel
        print(f"Reading Excel file: {file.filename}")
        df = pd.read_excel(file.file, sheet_name=0, header=0)
        tickets = parse_excel_to_tickets(df)
        
        print(f"Parsed {len(tickets)} tickets from Excel")
        
        # Create batch
        manager = get_batch_manager()
        batch = manager.create_batch(
            excel_filename=file.filename,
            total_passengers=len(tickets)
        )
        
        batch_id = batch["batch_id"]
        
        # Add all passengers to batch metadata
        for ticket in tickets:
            manager.add_passenger_to_batch(
                batch_id=batch_id,
                passenger_info={
                    "pax_name": ticket.pax_name,
                    "pnr": ticket.pnr,
                    "ticket_type": ticket.ticket_type
                }
            )
        
        # Start background generation
        background_tasks.add_task(generate_tickets_for_batch, batch_id, tickets)
        
        return UploadResponse(
            batch_id=batch_id,
            filename=file.filename,
            total_passengers=len(tickets),
            status="processing",
            message=f"Batch created successfully. Generating {len(tickets)} tickets..."
        )
        
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(
            status_code=500, 
            detail=f"Error processing file: {str(e)}"
        )


@router.get("/batches", response_model=BatchListResponse)
async def list_batches(
    page: int = Query(1, ge=1),
    limit: int = Query(10, ge=1, le=100)
):
    """
    Get paginated list of batches
    For the table view in Image 1
    """
    try:
        manager = get_batch_manager()
        result = manager.list_batches(page=page, limit=limit)
        
        return BatchListResponse(**result)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/batches/{batch_id}", response_model=BatchDetail)
async def get_batch_details(batch_id: str):
    """
    Get full batch details including all passengers
    For the detail view in Image 2
    """
    try:
        manager = get_batch_manager()
        batch_details = manager.get_batch_details(batch_id)
        
        if not batch_details:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        return BatchDetail(**batch_details)
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/batches/{batch_id}/status", response_model=BatchStatusResponse)
async def get_batch_status(batch_id: str):
    """
    Get real-time batch status
    Frontend polls this endpoint to show progress
    """
    try:
        manager = get_batch_manager()
        batch_details = manager.get_batch_details(batch_id)
        
        if not batch_details:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        total = batch_details["total_passengers"]
        generated = batch_details["generated"]
        failed = batch_details["failed"]
        pending = total - generated - failed
        
        progress = calculate_progress_percentage(generated, failed, total)
        is_complete = batch_details["status"] == "completed"
        
        return BatchStatusResponse(
            batch_id=batch_id,
            status=batch_details["status"],
            total_passengers=total,
            generated=generated,
            failed=failed,
            pending=pending,
            progress_percentage=progress,
            is_complete=is_complete,
            passengers=batch_details["passengers"]
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download/pdf/{batch_id}/{filename}")
async def download_single_pdf(batch_id: str, filename: str):
    """
    Download a single PDF file
    Triggered by download icon in UI
    """
    try:
        file_handler = get_file_handler()
        return file_handler.serve_pdf(batch_id, filename)
        
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="PDF file not found")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download/batch/{batch_id}")
async def download_batch_zip(batch_id: str):
    """
    Download all PDFs in a batch as ZIP
    Triggered by download icon for entire batch
    """
    try:
        manager = get_batch_manager()
        file_handler = get_file_handler()
        
        batch_details = manager.get_batch_details(batch_id)
        if not batch_details:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        return file_handler.serve_batch_zip(batch_id, batch_details)
        
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="No PDFs found in batch")
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.delete("/batches/{batch_id}")
async def delete_batch(batch_id: str):
    """
    Delete a batch and all its files
    Triggered by delete icon in UI
    """
    try:
        manager = get_batch_manager()
        
        # Check if batch exists
        batch_info = manager.get_batch_info(batch_id)
        if not batch_info:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        # Delete batch
        success = manager.delete_batch(batch_id)
        
        if success:
            return create_success_response(
                data={"batch_id": batch_id},
                message="Batch deleted successfully"
            )
        else:
            raise HTTPException(status_code=500, detail="Failed to delete batch")
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/stats", response_model=DashboardStats)
async def get_statistics():
    """
    Get overall statistics
    For the dashboard cards in Image 1
    """
    try:
        manager = get_batch_manager()
        stats = manager.get_statistics()
        
        return DashboardStats(**stats)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/system-status", response_model=SystemStatus)
async def check_system_status():
    """
    Check system status
    Useful for debugging and initial setup verification
    """
    try:
        converter = get_converter()
        manager = get_batch_manager()
        
        return SystemStatus(
            libreoffice_available=converter.is_available(),
            libreoffice_path=converter.libreoffice_path if converter.is_available() else None,
            output_directory_exists=manager.batches_dir.exists(),
            manifest_exists=manager.manifest_path.exists(),
            total_batches=len(manager._load_manifest()["batches"])
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))