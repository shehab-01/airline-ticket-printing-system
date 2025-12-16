from fastapi import APIRouter, File, UploadFile, HTTPException, BackgroundTasks, Query
from fastapi.responses import StreamingResponse, FileResponse
from api.model.models import (
    ResData,
    BatchListResponse,
    BatchInfo,
    BatchDetail,
    BatchStatusResponse,
    UploadResponse,
    DashboardStats,
    SystemStatus,
    PassengerInfo,
    create_success_response,
    calculate_progress_percentage
)
from api.utils.pdf_converter import get_converter
from api.utils.batch_manager import get_batch_manager
from api.utils.file_handler import get_file_handler
from api.utils.agency_manager import get_agency_manager  # ✅ NEW: Added agency manager
from api.utils.ticket_number_manager import get_ticket_number_manager  # ✅ NEW: Added ticket number manager
import pandas as pd
from dataclasses import dataclass
from typing import Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from datetime import datetime
import traceback


router = APIRouter(prefix="/api/v1/ticket")


# ============================================================================
# DATA CLASSES (Keep your existing Ticket dataclass)
# ============================================================================

@dataclass
class Ticket:
    no: int
    rsvn_cfmd: str
    ticket_type: str
    pax_name: str
    emd1: str
    pnr: str
    travel_agency: str
    ptn1_dep: str
    ptn1_dep_date: str
    ptn1_dep_time: str
    ptn1_arr: str
    ptn1_arr_date: str
    ptn1_arr_time: str
    ptn2_dep: Optional[str] = None
    ptn2_dep_date: Optional[str] = None
    ptn2_dep_time: Optional[str] = None
    ptn2_arr: Optional[str] = None
    ptn2_arr_date: Optional[str] = None
    ptn2_arr_time: Optional[str] = None


# ============================================================================
# HELPER FUNCTIONS (Keep your existing PPTX functions)
# ============================================================================

def replace_text_in_paragraph(paragraph, replacements):
    full_text = paragraph.text
    original_text = full_text
    
    for old_text, new_text in replacements.items():
        if old_text in full_text:
            full_text = full_text.replace(old_text, str(new_text))
    
    if full_text != original_text:
        # Preserve the first run's formatting
        if paragraph.runs:
            first_run = paragraph.runs[0]
            # Store formatting
            font_name = first_run.font.name
            font_size = first_run.font.size
            bold = first_run.font.bold
            italic = first_run.font.italic
            
            # Clear all runs
            for run in paragraph.runs:
                run.text = ""
            
            # Set new text with preserved formatting
            first_run.text = full_text
            if font_name:
                first_run.font.name = font_name
            if font_size:
                first_run.font.size = font_size
            if bold is not None:
                first_run.font.bold = bold
            if italic is not None:
                first_run.font.italic = italic
        else:
            paragraph.add_run(full_text)


# ============================================================================
# AIRPORT DATA MAPPING
# ============================================================================

AIRPORT_DB = {
    "ICN": {
        "city": "SEOUL",
        "airport": "Incheon International Airport Terminal - 1"
    },
    "DAC": {
        "city": "DHAKA",
        "airport": "Hazrat Shahjalal International Airport"
    },
    # Add more airports here as needed
    # "JFK": {"city": "NEW YORK", "airport": "John F. Kennedy International Airport"},
}

def get_airport_details(code: Optional[str]):
    
    if not code:
        return "", ""
    
    # Normalize code (uppercase, strip spaces)
    clean_code = str(code).strip().upper()
    
    data = AIRPORT_DB.get(clean_code)
    
    if data:
        return data["city"], data["airport"]
    
    # Fallback: if code unknown, just return the code to avoid crashing or empty space
    return clean_code, clean_code


def replace_text_in_shape(shape, replacements):
    
    # Handle groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            replace_text_in_shape(child_shape, replacements)
        return

    # Handle text frames
    if shape.has_text_frame:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            replace_text_in_paragraph(paragraph, replacements)
    
    # Handle tables
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replace_text_in_paragraph(paragraph, replacements)


def replace_image_in_slide(slide, placeholder_name: str, new_image_path: Path):
   
    if not new_image_path or not new_image_path.exists():
        print(f"  ⚠️  Image not found: {new_image_path}")
        return False
    
    for shape in slide.shapes:
        # Check if this is the placeholder image
        # Try both name and alt text (description)
        shape_name = shape.name if hasattr(shape, 'name') else ''
        
        # Check alt text (some versions store it differently)
        try:
            alt_text = shape._element.get('{http://schemas.openxmlformats.org/drawingml/2006/main}name', '')
        except:
            alt_text = ''
        
        # Match by name or alt text
        if placeholder_name.lower() in shape_name.lower() or placeholder_name.lower() in alt_text.lower():
            try:
                # Get position and size of placeholder
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Get the shape's index in the slide
                shape_idx = None
                for idx, s in enumerate(slide.shapes):
                    if s == shape:
                        shape_idx = idx
                        break
                
                # Delete the placeholder
                if shape_idx is not None:
                    sp = slide.shapes._spTree
                    sp.remove(shape._element)
                
                # Add new image at same position and size
                slide.shapes.add_picture(
                    str(new_image_path),
                    left, top,
                    width, height
                )
                
                print(f"  ✓ Logo replaced: {new_image_path.name}")
                return True
                
            except Exception as e:
                print(f"  ❌ Error replacing image: {e}")
                return False
    
    print(f"  Placeholder image '{placeholder_name}' not found in slide")
    return False


def generate_ppt_from_template(template_path, ticket, agency=None, ticket_number=None):

    prs = Presentation(template_path)
    
    current_date = datetime.now().strftime("%Y-%m-%d")
    ptn1_dep_city, ptn1_dep_airport = get_airport_details(ticket.ptn1_dep)
    ptn1_arr_city, ptn1_arr_airport = get_airport_details(ticket.ptn1_arr)
    
    replacements = {
        '{{date_now}}': current_date,
        '{{PAX_name}}': ticket.pax_name,
        '{{PNR_Reference}}': ticket.pnr,
        '{{Ticket_Number}}': ticket_number if ticket_number else '0000000000',  
        '{{EMD1}}': ticket.emd1 if ticket.emd1 else "0",
        '{{Ticket_Type}}' : ticket.ticket_type,
        '{{PTN1-Dep}}': ticket.ptn1_dep,
        '{{PTN1-Arr}}': ticket.ptn1_arr,
        '{{PTN1_Date}}': ticket.ptn1_dep_date,
        '{{PTN1_Time}}': ticket.ptn1_dep_time,
        '{{PTN1_Date.1}}': ticket.ptn1_arr_date,
        '{{PTN1_Time.1}}': ticket.ptn1_arr_time,

        '{{PTN1-Dep-City}}': ptn1_dep_city,
        '{{PTN1-Dep-Airport}}': ptn1_dep_airport,
        '{{PTN1-Arr-City}}': ptn1_arr_city,
        '{{PTN1-Arr-Airport}}': ptn1_arr_airport,
    }
    
    if ticket.ptn2_dep:
        ptn2_dep_city, ptn2_dep_airport = get_airport_details(ticket.ptn2_dep)
        ptn2_arr_city, ptn2_arr_airport = get_airport_details(ticket.ptn2_arr)

        replacements.update({
            '{{PTN2-Dep}}': ticket.ptn2_dep,
            '{{PTN2-Arr}}': ticket.ptn2_arr,
            '{{PTN2_Date}}': ticket.ptn2_dep_date,
            '{{PTN2_Time}}': ticket.ptn2_dep_time,
            '{{PTN2_Date.1}}': ticket.ptn2_arr_date,
            '{{PTN2_Time.1}}': ticket.ptn2_arr_time,

            '{{PTN2-Dep-City}}': ptn2_dep_city,
            '{{PTN2-Dep-Airport}}': ptn2_dep_airport,
            '{{PTN2-Arr-City}}': ptn2_arr_city,
            '{{PTN2-Arr-Airport}}': ptn2_arr_airport,
        })
    else:
        # Clear PTN2 placeholders if no return flight
        replacements.update({
            '{{PTN2-Dep}}': '',
            '{{PTN2-Arr}}': '',
            '{{PTN2_Date}}': '',
            '{{PTN2_Time}}': '',
            '{{PTN2_Date.1}}': '',
            '{{PTN2_Time.1}}': '',
            '{{PTN2-Dep-City}}': '',
            '{{PTN2-Dep-Airport}}': '',
            '{{PTN2-Arr-City}}': '',
            '{{PTN2-Arr-Airport}}': '',
        })
    
    #  Add agency information if found
    if agency:
        replacements.update({
            '{{agency_name}}': agency.get('agency_name', ''),
            '{{agency_owner}}': agency.get('agency_owner', ''),
            '{{agency_address}}': agency.get('agency_address', ''),
            '{{agency_email}}': agency.get('email', ''),
            '{{agency_phone}}': agency.get('telephone', '')
        })
        print(f" Using agency: {agency.get('agency_name')}")
    else:
        # Use defaults if agency not found
        replacements.update({
            '{{agency_name}}': ticket.travel_agency if ticket.travel_agency else '',
            '{{agency_owner}}': '',
            '{{agency_address}}': '',
            '{{agency_email}}': '',
            '{{agency_phone}}': ''
        })
        if ticket.travel_agency:
            print(f" Agency not found in database: {ticket.travel_agency}")
    
    # Replace text in all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
    
    # Replace agency logo if available
    if agency and agency.get('logo_path'):
        logo_path = Path(agency['logo_path'])
        if logo_path.exists():
            # Replace logo in all slides
            for slide in prs.slides:
                replace_image_in_slide(slide, 'agency_logo', logo_path)
        else:
            print(f"  ⚠️  Agency logo file not found: {logo_path}")
    
    return prs


def parse_excel_to_tickets(df: pd.DataFrame) -> list[Ticket]:

    df.columns = df.columns.str.replace(' ', '_')
    df = df.fillna("")
    
    tickets = []
    for index, row in df.iterrows():
        raw_emd = row.get('emd1_(Extra_RQ)', '')
        clean_emd = str(int(raw_emd)) if raw_emd != "" else ""

        ticket = Ticket(
            no=row['NO'],
            rsvn_cfmd=row['RSVN_cfmd'],
            ticket_type=row['ADT/LBR/CHD/INF'],
            pax_name=row['PAX_name'],
            emd1=clean_emd, 
            pnr=row['PNR_Reference'],
            travel_agency=row.get('Travel_Agency', ''),
            ptn1_dep=row['PTN1-Dep'],
            ptn1_dep_date=row['PTN1_Date'],
            ptn1_dep_time=row['PTN1_Time'],
            ptn1_arr=row['PTN1-Arr'],
            ptn1_arr_date=row['PTN1_Date.1'],  
            ptn1_arr_time=row['PTN1_Time.1'], 
            ptn2_dep=row.get('PTN2-Dep', ''),
            ptn2_dep_date=row.get('PTN2_Date', ''),
            ptn2_dep_time=row.get('PTN2_Time', ''),
            ptn2_arr=row.get('PTN2-Arr', ''),
            ptn2_arr_date=row.get('PTN2_Date.1', ''),
            ptn2_arr_time=row.get('PTN2_Time.1', '')
        )
        tickets.append(ticket)

    return tickets


# ============================================================================
# TICKET GENERATION LOGIC
# ============================================================================

def generate_tickets_for_batch(batch_id: str, tickets: list[Ticket]):
 
    manager = get_batch_manager()
    converter = get_converter()
    file_handler = get_file_handler()
    agency_manager = get_agency_manager()  # ✅ NEW: Get agency manager
    ticket_num_manager = get_ticket_number_manager()  # ✅ NEW: Get ticket number manager
    
    # template_path = Path("templates/ticket_template.pptx")
    batch_dir = manager.get_batch_dir(batch_id)
    
    # Update batch status to processing
    manager.update_batch_status(batch_id, "processing")
    
    print(f"\n{'='*60}")
    print(f"Starting batch generation: {batch_id}")
    print(f"Total tickets: {len(tickets)}")
    print(f"{'='*60}\n")
    
    for idx, ticket in enumerate(tickets, 1):
        print(f"[{idx}/{len(tickets)}] Processing: {ticket.pax_name} ({ticket.pnr})")
        
        try:
            if ticket.ptn2_dep:
                # Has return flight -> Use Round Trip Template
                template_path = Path("templates/ticket_template_roundtrip.pptx")
                print(f" Type: Round Trip")
            else:
                # No return flight -> Use One Way Template
                template_path = Path("templates/ticket_template_oneway.pptx")
                print(f" Type: One Way")

            # Check if template exists before proceeding
            if not template_path.exists():
                raise FileNotFoundError(f"Template not found: {template_path}")

            #  Look up agency from database
            agency = None
            if ticket.travel_agency:
                agency = agency_manager.find_by_name(ticket.travel_agency)
            
            #  Generate ticket number
            agency_name = agency.get('agency_name') if agency else ticket.travel_agency
            if not agency_name:
                agency_name = "Unknown"  # Fallback
            
            ticket_number = ticket_num_manager.generate_ticket_number(agency_name)
            
            #  Generate PPTX with agency info and ticket number
            prs = generate_ppt_from_template(template_path, ticket, agency, ticket_number)
            
            # Create safe filename
            safe_pax_name = ticket.pax_name.replace('/', '_').replace('\\', '_')
            safe_pnr = ticket.pnr.replace('/', '_')
            
            pptx_filename = f"{ticket_number}_{safe_pax_name}.pptx"
            pptx_path = batch_dir / pptx_filename
            
            # Save PPTX temporarily
            prs.save(str(pptx_path))
            print(f"  PPTX created: {pptx_filename}")
            
            # Convert PPTX to PDF
            pdf_path = converter.convert_pptx_to_pdf(pptx_path, batch_dir)
            pdf_filename = pdf_path.name
            print(f"  PDF created: {pdf_filename}")
            
            # Update passenger status to generated
            manager.update_passenger_status(
                batch_id=batch_id,
                pax_name=ticket.pax_name,
                pnr=ticket.pnr,
                status="generated",
                pdf_filename=pdf_filename
            )
            
            # Clean up PPTX file
            try:
                pptx_path.unlink()
                print(f"  Cleaned up PPTX")
            except:
                pass
            
            print(f"  Completed: {ticket.pax_name}\n")
            
        except Exception as e:
            error_msg = str(e)
            print(f"  Failed: {error_msg}\n")
            
            # Update passenger status to failed
            manager.update_passenger_status(
                batch_id=batch_id,
                pax_name=ticket.pax_name,
                pnr=ticket.pnr,
                status="failed",
                error=error_msg
            )
    
    # Mark batch as completed
    manager.update_batch_status(batch_id, "completed")
    
    # Get final counts
    batch_info = manager.get_batch_info(batch_id)
    print(f"\n{'='*60}")
    print(f"Batch {batch_id} completed!")
    print(f"Generated: {batch_info['generated']}")
    print(f"Failed: {batch_info['failed']}")
    print(f"{'='*60}\n")


# ============================================================================
# API ENDPOINTS
# ============================================================================

@router.post("/upload", response_model=UploadResponse)
async def upload_excel(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...)
):

    # Validate file type
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(
            status_code=400, 
            detail="File must be an Excel file (.xlsx or .xls)"
        )
    
    # Check LibreOffice availability
    converter = get_converter()
    if not converter.is_available():
        raise HTTPException(
            status_code=500,
            detail="LibreOffice not found. Please install LibreOffice from https://www.libreoffice.org/download/"
        )
    
    try:
        # Parse Excel
        print(f"Reading Excel file: {file.filename}")
        df = pd.read_excel(file.file, sheet_name=0, header=0)
        tickets = parse_excel_to_tickets(df)
        
        print(f"Parsed {len(tickets)} tickets from Excel")
        
        # Create batch
        manager = get_batch_manager()
        batch = manager.create_batch(
            excel_filename=file.filename,
            total_passengers=len(tickets)
        )
        
        batch_id = batch["batch_id"]
        
        # Add all passengers to batch metadata
        for ticket in tickets:
            manager.add_passenger_to_batch(
                batch_id=batch_id,
                passenger_info={
                    "pax_name": ticket.pax_name,
                    "pnr": ticket.pnr,
                    "ticket_type": ticket.ticket_type
                }
            )
        
        # Start background generation
        background_tasks.add_task(generate_tickets_for_batch, batch_id, tickets)
        
        return UploadResponse(
            batch_id=batch_id,
            filename=file.filename,
            total_passengers=len(tickets),
            status="processing",
            message=f"Batch created successfully. Generating {len(tickets)} tickets..."
        )
        
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(
            status_code=500, 
            detail=f"Error processing file: {str(e)}"
        )


@router.get("/batches", response_model=BatchListResponse)
async def list_batches(
    page: int = Query(1, ge=1),
    limit: int = Query(10, ge=1, le=100)
):

    try:
        manager = get_batch_manager()
        result = manager.list_batches(page=page, limit=limit)
        
        return BatchListResponse(**result)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/batches/{batch_id}", response_model=BatchDetail)
async def get_batch_details(batch_id: str):
    
    try:
        manager = get_batch_manager()
        batch_details = manager.get_batch_details(batch_id)
        
        if not batch_details:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        return BatchDetail(**batch_details)
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/batches/{batch_id}/status", response_model=BatchStatusResponse)
async def get_batch_status(batch_id: str):
    
    try:
        manager = get_batch_manager()
        batch_details = manager.get_batch_details(batch_id)
        
        if not batch_details:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        total = batch_details["total_passengers"]
        generated = batch_details["generated"]
        failed = batch_details["failed"]
        pending = total - generated - failed
        
        progress = calculate_progress_percentage(generated, failed, total)
        is_complete = batch_details["status"] == "completed"
        
        return BatchStatusResponse(
            batch_id=batch_id,
            status=batch_details["status"],
            total_passengers=total,
            generated=generated,
            failed=failed,
            pending=pending,
            progress_percentage=progress,
            is_complete=is_complete,
            passengers=batch_details["passengers"]
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download/pdf/{batch_id}/{filename}")
async def download_single_pdf(batch_id: str, filename: str):
    
    try:
        file_handler = get_file_handler()
        return file_handler.serve_pdf(batch_id, filename)
        
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="PDF file not found")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download/batch/{batch_id}")
async def download_batch_zip(batch_id: str):
    
    try:
        manager = get_batch_manager()
        file_handler = get_file_handler()
        
        batch_details = manager.get_batch_details(batch_id)
        if not batch_details:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        return file_handler.serve_batch_zip(batch_id, batch_details)
        
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="No PDFs found in batch")
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.delete("/batches/{batch_id}")
async def delete_batch(batch_id: str):
    
    try:
        manager = get_batch_manager()
        
        # Check if batch exists
        batch_info = manager.get_batch_info(batch_id)
        if not batch_info:
            raise HTTPException(status_code=404, detail="Batch not found")
        
        # Delete batch
        success = manager.delete_batch(batch_id)
        
        if success:
            return create_success_response(
                data={"batch_id": batch_id},
                message="Batch deleted successfully"
            )
        else:
            raise HTTPException(status_code=500, detail="Failed to delete batch")
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/stats", response_model=DashboardStats)
async def get_statistics():
    
    try:
        manager = get_batch_manager()
        stats = manager.get_statistics()
        
        return DashboardStats(**stats)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/system-status", response_model=SystemStatus)
async def check_system_status():
    
    try:
        converter = get_converter()
        manager = get_batch_manager()
        
        return SystemStatus(
            libreoffice_available=converter.is_available(),
            libreoffice_path=converter.libreoffice_path if converter.is_available() else None,
            output_directory_exists=manager.batches_dir.exists(),
            manifest_exists=manager.manifest_path.exists(),
            total_batches=len(manager._load_manifest()["batches"])
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))