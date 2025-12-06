from fastapi import FastAPI, APIRouter, File, Form, HTTPException, UploadFile, Depends
from fastapi.responses import StreamingResponse
from api.model.models import ResData
import pandas as pd
import re
from dataclasses import dataclass, asdict
from typing import Optional
from docxtpl import DocxTemplate
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import os
from pathlib import Path
from datetime import datetime
import zipfile


router = APIRouter(prefix="/api/v1/ticket")

@dataclass
class Ticket:
    no: int
    rsvn_cfmd: str
    ticket_type: str
    pax_name: str
    emd1: str
    pnr: str
    travel_agency: str
    ptn1_dep: str
    ptn1_dep_date: str
    ptn1_dep_time: str
    ptn1_arr: str
    ptn1_arr_date: str
    ptn1_arr_time: str
    ptn2_dep: Optional[str] = None
    ptn2_dep_date: Optional[str] = None
    ptn2_dep_time: Optional[str] = None
    ptn2_arr: Optional[str] = None
    ptn2_arr_date: Optional[str] = None
    ptn2_arr_time: Optional[str] = None


def replace_text_in_paragraph(paragraph, replacements):
    """Replace text in a paragraph by combining all runs first"""
    # Get the full text
    full_text = paragraph.text
    
    # Check if any placeholder exists in the full text
    for old_text, new_text in replacements.items():
        if old_text in full_text:
            full_text = full_text.replace(old_text, str(new_text))
    
    # If text was changed, update the paragraph
    if full_text != paragraph.text:
        # Clear all runs
        for run in paragraph.runs:
            run.text = ""
        
        # Set the new text in the first run
        if paragraph.runs:
            paragraph.runs[0].text = full_text
        else:
            paragraph.add_run(full_text)


# def replace_text_in_shape(shape, replacements):
#     """Replace text in a shape"""
#     if shape.has_text_frame:
#         text_frame = shape.text_frame
#         for paragraph in text_frame.paragraphs:
#             replace_text_in_paragraph(paragraph, replacements)
    
#     # Handle tables
#     if shape.has_table:
#         for row in shape.table.rows:
#             for cell in row.cells:
#                 for paragraph in cell.text_frame.paragraphs:
#                     replace_text_in_paragraph(paragraph, replacements)

def replace_text_in_shape(shape, replacements):
    """Replace text in a shape, handling Groups, Tables, and TextFrames recursively"""
    
    # 1. HANDLE GROUPS (This is likely the missing piece)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            replace_text_in_shape(child_shape, replacements)
        return

    # 2. Handle Text Frames (Standard Text Boxes)
    if shape.has_text_frame:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            replace_text_in_paragraph(paragraph, replacements)
    
    # 3. Handle Tables
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replace_text_in_paragraph(paragraph, replacements)

def generate_ppt_from_template(template_path, ticket):
    """Generate PowerPoint from template with ticket data"""
    prs = Presentation(template_path)
    
    # Get current date for the ticket issued date
    current_date = datetime.now().strftime("%Y-%m-%d")
    
    # Create replacements dictionary matching your template
    replacements = {
        '{{date_now}}': current_date,
        '{{PAX_name}}': ticket.pax_name,
        '{{PNR_Reference}}': ticket.pnr,
        '{{PTN1-Dep}}': ticket.ptn1_dep,
        '{{PTN1-Arr}}': ticket.ptn1_arr,
        '{{PTN1_Date}}': ticket.ptn1_dep_date,
        '{{PTN1_Time}}': ticket.ptn1_dep_time,
        '{{PTN1_Date.1}}': ticket.ptn1_arr_date,
        '{{PTN1_Time.1}}': ticket.ptn1_arr_time,
    }
    
    # Add PTN2 if exists (return flight)
    if ticket.ptn2_dep:
        replacements.update({
            '{{PTN2-Dep}}': ticket.ptn2_dep,
            '{{PTN2-Arr}}': ticket.ptn2_arr,
            '{{PTN2_Date}}': ticket.ptn2_dep_date,
            '{{PTN2_Time}}': ticket.ptn2_dep_time,
            '{{PTN2_Date.1}}': ticket.ptn2_arr_date,
            '{{PTN2_Time.1}}': ticket.ptn2_arr_time,
        })
    
    # Replace text in all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
    
    return prs


@router.post("/generate-ppt", response_model=ResData)
async def generate_ppt(file: UploadFile = File(...)):
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(status_code=400, detail="File must be an Excel file (.xlsx or .xls)")
    
    try:
        print("Reading the file")
        df = pd.read_excel(file.file, sheet_name=0, header=0)
        df.columns = df.columns.str.replace(' ', '_')
        df = df.fillna("")
        
        # Create tickets
        tickets = []
        for index, row in df.iterrows():
            raw_emd = row.get('emd1_(Extra_RQ)', '')
            clean_emd = str(int(raw_emd)) if raw_emd != "" else ""

            ticket = Ticket(
                no=row['NO'],
                rsvn_cfmd=row['RSVN_cfmd'],
                ticket_type=row['ADT/LBR/CHD/INF'],
                pax_name=row['PAX_name'],
                emd1=clean_emd, 
                pnr=row['PNR_Reference'],
                travel_agency=row['Travel_Agency',''],
                ptn1_dep=row['PTN1-Dep'],
                ptn1_dep_date=row['PTN1_Date'],
                ptn1_dep_time=row['PTN1_Time'],
                ptn1_arr=row['PTN1-Arr'],
                ptn1_arr_date=row['PTN1_Date.1'],  
                ptn1_arr_time=row['PTN1_Time.1'], 
                ptn2_dep=row.get('PTN2-Dep', ''),
                ptn2_dep_date=row.get('PTN2_Date', ''),
                ptn2_dep_time=row.get('PTN2_Time', ''),
                ptn2_arr=row.get('PTN2-Arr', ''),
                ptn2_arr_date=row.get('PTN2_Date.1', ''),
                ptn2_arr_time=row.get('PTN2_Time.1', '')
            )
            tickets.append(ticket)
        
        print(f"Created {len(tickets)} ticket objects")
        print(tickets)

        # Create tickets directory if it doesn't exist
        tickets_dir = Path("tickets")
        tickets_dir.mkdir(exist_ok=True)
        
        # Load the template
        template_path = "templates/ticket_template.pptx"
        
        # Generate presentations for all tickets
        generated_files = []
        for ticket in tickets:
            prs = generate_ppt_from_template(template_path, ticket)
            
            # Create a safe filename
            safe_pax_name = ticket.pax_name.replace('/', '_').replace('\\', '_')
            filename = f"ticket_{safe_pax_name}_{ticket.pnr.replace('/', '_')}.pptx"
            filepath = tickets_dir / filename
            
            # Save to the tickets directory
            prs.save(str(filepath))
            generated_files.append(filename)
            
            print(f"Generated: {filename}")
        
        return ResData(
            data={'generated_files': generated_files, 'count': len(generated_files)}, 
            msg=f"Successfully generated {len(generated_files)} PowerPoint presentations"
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")


@router.post("/generate-both")
async def generate_both(file: UploadFile = File(...)):
    """Generate both Word and PowerPoint documents"""
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(status_code=400, detail="File must be an Excel file (.xlsx or .xls)")
    
    try:
        print("Reading the file")
        df = pd.read_excel(file.file, sheet_name=0, header=0)
        df.columns = df.columns.str.replace(' ', '_')
        df = df.fillna("")
        
        # Create tickets
        tickets = []
        for index, row in df.iterrows():
            raw_emd = row.get('emd1_(Extra_RQ)', '')
            clean_emd = str(int(raw_emd)) if raw_emd != "" else ""

            ticket = Ticket(
                no=row['NO'],
                rsvn_cfmd=row['RSVN_cfmd'],
                ticket_type=row['ADT/LBR/CHD/INF'],
                pax_name=row['PAX_name'],
                emd1=clean_emd, 
                pnr=row['PNR_Reference'],
                travel_agency=row['Travel_Agency'],
                ptn1_dep=row['PTN1-Dep'],
                ptn1_dep_date=row['PTN1_Date'],
                ptn1_dep_time=row['PTN1_Time'],
                ptn1_arr=row['PTN1-Arr'],
                ptn1_arr_date=row['PTN1_Date.1'],  
                ptn1_arr_time=row['PTN1_Time.1'], 
                ptn2_dep=row.get('PTN2-Dep', ''),
                ptn2_dep_date=row.get('PTN2_Date', ''),
                ptn2_dep_time=row.get('PTN2_Time', ''),
                ptn2_arr=row.get('PTN2-Arr', ''),
                ptn2_arr_date=row.get('PTN2_Date.1', ''),
                ptn2_arr_time=row.get('PTN2_Time.1', '')
            )
            tickets.append(ticket)
        
        print(f"Created {len(tickets)} ticket objects")

        # Create tickets directory
        tickets_dir = Path("tickets")
        tickets_dir.mkdir(exist_ok=True)
        
        # Templates
        word_template = "Victory_Travels.docx"
        ppt_template = "Victory_Travels.pptx"
        
        # Create ZIP file
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for ticket in tickets:
                safe_pax_name = ticket.pax_name.replace('/', '_').replace('\\', '_')
                safe_pnr = ticket.pnr.replace('/', '_')
                context = asdict(ticket)
                
                # Add current date to context
                context['date_now'] = datetime.now().strftime("%Y-%m-%d")
                
                # Generate Word document
                doc = DocxTemplate(word_template)
                doc.render(context)
                
                doc_filename = f"ticket_{safe_pax_name}_{safe_pnr}.docx"
                doc_filepath = tickets_dir / doc_filename
                doc.save(str(doc_filepath))
                
                doc_buffer = io.BytesIO()
                doc.save(doc_buffer)
                doc_buffer.seek(0)
                zip_file.writestr(f"word/{doc_filename}", doc_buffer.getvalue())
                
                # Generate PowerPoint
                prs = generate_ppt_from_template(ppt_template, ticket)
                
                ppt_filename = f"ticket_{safe_pax_name}_{safe_pnr}.pptx"
                ppt_filepath = tickets_dir / ppt_filename
                prs.save(str(ppt_filepath))
                
                # Read the saved file for ZIP
                with open(str(ppt_filepath), 'rb') as f:
                    zip_file.writestr(f"ppt/{ppt_filename}", f.read())
                
                print(f"Generated: {doc_filename} and {ppt_filename}")
        
        zip_buffer.seek(0)
        
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=all_tickets.zip"}
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")